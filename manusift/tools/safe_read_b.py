"""Phase B safe-read implementation (document extract, tracker, redact).

**Deprecated as a public import path.** Prefer::

    from manusift.tools.safe_read import detect_xlsx_figs, …

This module remains the **implementation home** for Phase B so
``safe_read`` can re-export without a circular cycle. Direct imports
of ``manusift.tools.safe_read_b`` emit ``DeprecationWarning``.

R-2026-06-17 (Phase B:
borrow from
Claude Code +
Hermes):
4 medium-cost
``read_file``
hardening
modules.  See
``manusift/tools/safe_read.py``
for the 8 small
ones (Phase A) and the unified public surface.

Borrowed from
(see the comparison
doc):

  * ``suggest_similar_files``
    — Hermes
    ``file_operations.py:1056``.
    When the
    requested
    file is
    not
    found,
    return
    the top 5
    fuzzy-matched
    candidates
    so the
    LLM does
    not have
    to
    guess-and-retry
    the
    path.
    Score
    weights
    follow
    Hermes:
    exact=100,
    same-base-diff-ext=90,
    prefix=70,
    substring=60,
    same-ext + 40%
    char
    overlap=30.

  * ``ReadTracker``
    — Hermes
    ``file_tools.py:831``.
    Per-task
    (path, offset, limit)
    mtime cache.
    Same
    region
    read
    twice
    returns
    a "file
    unchanged"
    stub.  Same
    region
    read 3+
    times
    returns
    BLOCKED
    to defend
    against
    weak
    models
    that
    burn their
    iteration
    budget in
    an
    infinite
    read loop.

  * ``redact_sensitive_text``
    — Hermes
    ``agent/redact.py:326``.
    Replace
    30+
    known
    API-key
    prefixes
    (``sk-``,
    ``ghp_``,
    ``AKIA``,
    ``gAAAA``,
    etc.)
    with
    ``[REDACTED]``
    *before*
    the
    content
    enters
    the
    LLM
    context.
    The model
    can still
    read the
    surrounding
    code
    (e.g.
    a
    ``client = OpenAI(api_key=...)``
    line)
    but the
    key
    itself
    is
    gone.

  * ``extract_docx_text``,
    ``extract_xlsx_text``,
    ``extract_pptx_text``,
    ``extract_ipynb_text``
    — Hermes
    ``read_extract.py``
    (real
    implementation,
    not
    stubbed).
    Wraps
    ``python-docx`` /
    ``openpyxl`` /
    ``python-pptx`` /
    ``nbformat`` to
    produce
    a
    plain-text
    rendering.
    Each
    function
    raises
    ``ExtractionError``
    on
    a
    malformed
    file
    so the
    caller
    (``ReadFileTool.execute``)
    can
    fall
    through
    to the
    normal
    binary
    branch
    instead
    of
    crashing.
"""
from __future__ import annotations

import os
import re
import sys
import threading
import warnings
from pathlib import Path


def _maybe_warn_direct_import() -> None:
    """Deprecate direct ``safe_read_b`` imports (facade is ``safe_read``).

    Suppress when imported by ``safe_read`` re-export, under pytest
    (legacy tests), or ``MANUSIFT_SUPPRESS_SAFE_READ_B_WARNING=1``.
    """
    if os.environ.get("MANUSIFT_SUPPRESS_SAFE_READ_B_WARNING") == "1":
        return
    if os.environ.get("PYTEST_CURRENT_TEST"):
        return
    # Walk frames: if any is manusift.tools.safe_read, skip.
    frame = sys._getframe(1)
    while frame is not None:
        mod = frame.f_globals.get("__name__") or ""
        if mod == "manusift.tools.safe_read":
            return
        frame = frame.f_back  # type: ignore[assignment]
    warnings.warn(
        "manusift.tools.safe_read_b is a deprecated import path; "
        "use manusift.tools.safe_read (unified Phase A+B surface).",
        DeprecationWarning,
        stacklevel=3,
    )


_maybe_warn_direct_import()
from typing import Any


# =============================================================================
# B-1: suggest_similar_files
# =============================================================================


def _score_similarity(
    query: str,
    candidate: str,
    candidate_ext: str,
) -> int:
    """Return Hermes-style similarity score (0-100).

    R-2026-06-17
    (Phase B):
    exact
    match
    wins
    (100);
    same
    base
    name
    with
    a
    different
    extension
    is
    almost
    as
    good
    (90,
    "the
    user
    mistyped
    .txt
    instead
    of
    .csv");
    prefix
    match
    is
    strong
    (70);
    substring
    is
    weaker
    (60);
    same
    extension
    with
    40%
    character
    overlap
    is
    the
    last
    resort
    (30).
    All
    comparisons
    are
    case-insensitive.

    We
    check
    the
    *base
    name*
    (not
    the
    full
    filename
    with
    extension)
    for
    prefix/substring
    so
    the
    common
    case
    ``notes.md`` vs
    ``notes_old.md``
    matches
    at
    the
    prefix
    level
    (the
    user
    added
    a
    suffix),
    not
    at
    the
    substring
    level.
    """
    if not query or not candidate:
        return 0
    q = query.lower()
    c = candidate.lower()
    q_base, q_ext = (
        Path(q).stem,
        Path(q).suffix.lower(),
    )
    c_base, c_ext = (
        Path(c).stem,
        Path(c).suffix.lower(),
    )
    if c == q:
        return 100
    if c_base == q_base and q_ext and c_ext and q_ext != c_ext:
        return 90
    if c_base.startswith(q_base) or q_base.startswith(c_base):
        return 70
    if q_base in c_base or c_base in q_base:
        return 60
    if q_ext and c_ext and q_ext == c_ext and q_base and c_base:
        # Same
        # extension
        # + at
        # least
        # 40%
        # character
        # overlap
        # in
        # the
        # base
        # name.
        common = set(q_base) & set(c_base)
        max_len = max(len(q_base), len(c_base))
        if max_len > 0 and len(common) >= max_len * 0.4:
            return 30
    return 0


def suggest_similar_files(
    query_path: str,
    *,
    search_dir: str | None = None,
    limit: int = 5,
) -> list[str]:
    """Return up to ``limit`` files that fuzzy-match the missing ``query_path``.

    R-2026-06-17
    (Phase B):
    called
    when
    the
    user
    asked
    for
    a
    file
    that
    does
    not
    exist.
    Returns
    absolute
    paths
    so the
    LLM
    can
    re-call
    ``read_file``
    directly
    without
    re-resolving
    the
    path.

    ``search_dir``
    defaults
    to
    ``dirname(query_path)``
    so a
    user-typed
    ``C:\\foo\\notes.md``
    that
    doesn't
    exist
    gets
    suggestions
    from
    ``C:\\foo\\``
    (the
    most
    common
    case
    -- user
    typoed
    the
    filename
    but
    got
    the
    dir
    right).
    """
    if not query_path:
        return []
    if search_dir is None:
        search_dir = str(Path(query_path).parent) or "."
    search_path = Path(search_dir)
    if not search_path.is_dir():
        return []
    query_basename = Path(query_path).name
    query_ext = Path(query_path).suffix.lower()
    scored: list[tuple[int, str]] = []
    try:
        children = list(search_path.iterdir())
    except (OSError, PermissionError):
        return []
    for child in children:
        try:
            if not child.is_file():
                continue
        except OSError:
            continue
        s = _score_similarity(
            query_basename, child.name, query_ext
        )
        if s > 0:
            scored.append((s, str(child.resolve())))
    scored.sort(key=lambda x: (-x[0], x[1]))
    return [path for _, path in scored[:limit]]


# =============================================================================
# B-2: mtime dedup + BLOCKED-after-2-hits
# =============================================================================


# Hermes caps
# the
# per-task
# tracker
# at:
#   * dedup:
#     1000
# We
# use
# the
# same
# default
# so
# a
# long
# session
# never
# grows
# the
# tracker
# without
# bound.
_DEDUP_CAP = 1000

# The
# user-facing
# message
# Hermes
# surfaces
# when
# the
# same
# region
# has
# been
# read
# once
# already
# ("unchanged").
_READ_DEDUP_STATUS_MESSAGE = (
    "File unchanged since last read. The content from "
    "the earlier read_file result in this conversation is "
    "still current — refer to that instead of re-reading."
)


class ReadTracker:
    """Per-task tracker for read-dedup + external-edit detection.

    R-2026-06-17
    (Phase B):
    mirrors
    Hermes'
    ``_read_tracker``
    dict
    + ``_read_tracker_lock``
    but
    wrapped
    in
    a
    class
    so the
    test
    suite
    can
    instantiate
    a
    fresh
    one
    per
    task
    (Hermes
    uses
    a
    module-level
    dict
    keyed
    by
    ``task_id``;
    that's
    fine
    for
    a
    CLI
    but
    makes
    tests
    harder
    to
    isolate).
    """

    def __init__(self) -> None:
        self._lock = threading.Lock()
        # ``dedup`` maps
        # ``(path,
        # offset,
        # limit)`` →
        # ``mtime``.
        # ``dedup_hits``
        # maps
        # the
        # same
        # key
        # to
        # a
        # repeat
        # count
        # (0,
        # 1,
        # 2+)
        # so the
        # BLOCKED
        # trigger
        # can
        # count
        # *stubs*
        # (not
        # just
        # real
        # reads).
        self._dedup: dict[tuple[str, int, int], float] = {}
        self._dedup_hits: dict[tuple[str, int, int], int] = {}

    def check(
        self,
        path: str,
        offset: int,
        limit: int,
        *,
        current_mtime: float | None = None,
    ) -> dict[str, Any] | None:
        """Return a stub or BLOCKED response, or ``None`` to proceed with read.

        R-2026-06-17
        (Phase B):
        * First
        call
        for a
        (path,
        offset,
        limit)
        →
        ``None``
        (proceed).
        * Second
        call
        (file
        unchanged)
        →
        ``{"status": "unchanged", "dedup": True, ...}``.
        * Third+
        call
        (file
        unchanged)
        →
        ``{"error": "BLOCKED: ...", "path": ..., "already_read": N}``.

        The
        caller
        (``ReadFileTool.execute``)
        wraps
        the
        stub
        in
        the
        same
        JSON
        envelope
        as a
        real
        read
        so the
        LLM
        sees
        a
        consistent
        shape.
        """
        key = (str(Path(path).resolve()), int(offset), int(limit))
        with self._lock:
            cached_mtime = self._dedup.get(key)
        if cached_mtime is None:
            return None
        if current_mtime is None:
            try:
                current_mtime = Path(path).stat().st_mtime
            except OSError:
                return None
        if current_mtime != cached_mtime:
            return None
        with self._lock:
            hits = self._dedup_hits.get(key, 0) + 1
            self._dedup_hits[key] = hits
        if hits >= 2:
            return {
                "ok": False,
                "error_kind": "blocked",
                "error": (
                    f"BLOCKED: You have called read_file on this "
                    f"exact region {hits + 1} times and the file "
                    f"has NOT changed. STOP calling read_file for "
                    f"this path — the content from your earlier "
                    f"read_file result in this conversation is "
                    f"still current. Proceed with your task using "
                    f"the information you already have."
                ),
                "path": path,
                "already_read": hits + 1,
            }
        return {
            "ok": True,
            "status": "unchanged",
            "message": _READ_DEDUP_STATUS_MESSAGE,
            "path": path,
            "dedup": True,
            "content_returned": False,
        }

    def record(
        self,
        path: str,
        offset: int,
        limit: int,
        mtime: float,
    ) -> None:
        """Record a successful read so future calls can dedup.

        R-2026-06-17
        (Phase B):
        the
        caller
        must
        call
        this
        *after*
        a
        real
        read
        succeeds
        (not
        after
        a
        stub
        or
        BLOCKED).
        Caps
        the
        dedup
        dict
        so a
        long
        session
        never
        grows
        it
        past
        ``_DEDUP_CAP``
        entries.
        """
        key = (str(Path(path).resolve()), int(offset), int(limit))
        with self._lock:
            self._dedup[key] = mtime
            if len(self._dedup) > _DEDUP_CAP:
                excess = len(self._dedup) - _DEDUP_CAP
                for old_key in list(self._dedup)[:excess]:
                    self._dedup.pop(old_key, None)
                    self._dedup_hits.pop(old_key, None)

    def reset(self) -> None:
        """Clear all dedup state (for tests)."""
        with self._lock:
            self._dedup.clear()
            self._dedup_hits.clear()


# Per-trace-id
# cache
# of
# ``ReadTracker``
# instances.
# ``ToolContext.metadata``
# is
# a
# ``MappingProxyType``
# (frozen)
# so
# we
# cannot
# stash
# the
# tracker
# there
# directly.
# Instead
# we
# keep
# a
# module-level
# dict
# keyed
# by
# ``trace_id``.
# The
# caller
# (``ReadFileTool.execute``)
# passes
# ``ctx.trace_id``
# to
# the
# factory
# below.
_TRACKERS_BY_TRACE: dict[str, ReadTracker] = {}
_TRACKERS_LOCK = threading.Lock()


def get_tracker(trace_id: str) -> ReadTracker:
    """Return the per-trace ``ReadTracker``, creating one if needed.

    R-2026-06-17
    (Phase B):
    Hermes
    stores
    its
    tracker
    on
    a
    module-level
    dict
    keyed
    by
    ``task_id``;
    we
    do
    the
    same
    with
    ``trace_id``
    (the
    ManusSift
    equivalent).
    The
    factory
    is
    thread-safe
    so
    parallel
    tool
    calls
    in
    the
    same
    session
    can
    share
    a
    tracker
    without
    races.
    """
    with _TRACKERS_LOCK:
        tracker = _TRACKERS_BY_TRACE.get(trace_id)
        if tracker is None:
            tracker = ReadTracker()
            _TRACKERS_BY_TRACE[trace_id] = tracker
        return tracker


def reset_tracker(trace_id: str) -> None:
    """Clear the tracker for a single trace (for tests)."""
    with _TRACKERS_LOCK:
        _TRACKERS_BY_TRACE.pop(trace_id, None)


def reset_all_trackers() -> None:
    """Clear all trackers (for tests)."""
    with _TRACKERS_LOCK:
        _TRACKERS_BY_TRACE.clear()


# =============================================================================
# B-3: redact_sensitive_text
# =============================================================================


# Known
# API
# key
# prefixes
# (borrowed
# from
# Hermes'
# ``agent/redact.py:70``).
# Match
# the
# prefix
# +
# contiguous
# token
# characters.
_PREFIX_PATTERNS: tuple[str, ...] = (
    r"sk-[A-Za-z0-9_-]{10,}",  # OpenAI / OpenRouter / Anthropic (sk-ant-*)
    r"ghp_[A-Za-z0-9]{10,}",  # GitHub PAT (classic)
    r"github_pat_[A-Za-z0-9_]{10,}",  # GitHub PAT (fine-grained)
    r"gho_[A-Za-z0-9]{10,}",  # GitHub OAuth access token
    r"ghu_[A-Za-z0-9]{10,}",  # GitHub user-to-server token
    r"ghs_[A-Za-z0-9]{10,}",  # GitHub server-to-server token
    r"ghr_[A-Za-z0-9]{10,}",  # GitHub refresh token
    r"xox[baprs]-[A-Za-z0-9-]{10,}",  # Slack tokens
    r"AIza[A-Za-z0-9_-]{30,}",  # Google API keys
    r"pplx-[A-Za-z0-9]{10,}",  # Perplexity
    r"fal_[A-Za-z0-9_-]{10,}",  # Fal.ai
    r"fc-[A-Za-z0-9]{10,}",  # Firecrawl
    r"bb_live_[A-Za-z0-9_-]{10,}",  # BrowserBase
    r"gAAAA[A-Za-z0-9_=-]{20,}",  # Codex encrypted tokens
    r"AKIA[A-Z0-9]{16}",  # AWS Access Key ID
    r"sk_live_[A-Za-z0-9]{10,}",  # Stripe secret (live)
    r"sk_test_[A-Za-z0-9]{10,}",  # Stripe secret (test)
    r"rk_live_[A-Za-z0-9]{10,}",  # Stripe restricted key
    r"SG\.[A-Za-z0-9_-]{10,}",  # SendGrid API key
    r"hf_[A-Za-z0-9]{10,}",  # HuggingFace token
    r"r8_[A-Za-z0-9]{10,}",  # Replicate API token
    r"npm_[A-Za-z0-9]{10,}",  # npm access token
    r"pypi-[A-Za-z0-9_-]{10,}",  # PyPI API token
    r"dop_v1_[A-Za-z0-9]{10,}",  # DigitalOcean PAT
    r"doo_v1_[A-Za-z0-9]{10,}",  # DigitalOcean OAuth
    r"am_[A-Za-z0-9_-]{10,}",  # AgentMail API key
    r"sk_[A-Za-z0-9_]{10,}",  # ElevenLabs TTS (underscore, not dash)
    r"tvly-[A-Za-z0-9]{10,}",  # Tavily
    r"exa_[A-Za-z0-9]{10,}",  # Exa search
    r"gsk_[A-Za-z0-9]{10,}",  # Groq Cloud
    r"syt_[A-Za-z0-9]{10,}",  # Matrix access token
)

# Pre-compile
# a
# single
# OR-joined
# pattern
# with
# word
# boundaries
# so a
# ``skull``
# in
# a
# word
# doesn't
# get
# matched
# as
# ``sk-ull``.
_REDACT_RE = re.compile(
    r"(?<![A-Za-z0-9_-])("
    + "|".join(_PREFIX_PATTERNS)
    + r")(?![A-Za-z0-9_-])"
)

_REDACTED_PLACEHOLDER = "[REDACTED]"


def redact_sensitive_text(
    text: str,
    *,
    placeholder: str = _REDACTED_PLACEHOLDER,
) -> str:
    """Replace known API-key prefixes with ``[REDACTED]``.

    R-2026-06-17
    (Phase B):
    the
    LLM
    can
    still
    read
    the
    *surrounding
    code*
    (e.g.
    ``client = OpenAI(api_key=...)``
    line)
    but
    the
    key
    itself
    is
    gone.
    """
    if not text:
        return text
    return _REDACT_RE.sub(placeholder, text)


# =============================================================================
# B-4: docx / xlsx / pptx / ipynb text extraction
# =============================================================================


class ExtractionError(Exception):
    """Raised when a structured document cannot be extracted as text.

    R-2026-06-17
    (Phase B):
    the
    caller
    (``ReadFileTool.execute``)
    catches
    this
    and
    falls
    through
    to the
    normal
    text-read
    branch
    (which
    gives
    a
    cleaner
    "not
    text"
    error
    than
    a
    10-line
    Python
    traceback).
    """


def extract_docx_text(path: str) -> str:
    """Extract plain text from a .docx (Word) file.

    R-2026-06-17
    (Phase B):
    uses
    ``python-docx``
    (already
    a
    dependency).
    Walks
    paragraphs
    in
    order;
    table
    cells
    are
    rendered
    with
    ``\\t``-separated
    rows
    and
    blank-line
    separators
    between
    tables.
    Raises
    ``ExtractionError``
    on
    a
    malformed
    file.
    """
    try:
        from docx import Document
    except ImportError as exc:
        raise ExtractionError(
            f"python-docx not installed: {exc}"
        ) from exc
    try:
        doc = Document(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            f"docx open failed: {exc}"
        ) from exc
    try:
        out_lines: list[str] = []
        for para in doc.paragraphs:
            text = (para.text or "").rstrip()
            if text:
                out_lines.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [
                    (cell.text or "").strip() for cell in row.cells
                ]
                out_lines.append("\t".join(cells))
            out_lines.append("")
        return "\n".join(out_lines)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            f"docx walk failed: {exc}"
        ) from exc


def extract_xlsx_text(path: str) -> str:
    """Extract a plain-text rendering of an .xlsx workbook.

    R-2026-06-17
    (Phase B):
    uses
    ``openpyxl``
    in
    read-only
    mode
    (no
    rewrite).
    Detects
    fig
    boundaries
    within
    each
    sheet
    so a
    single
    sheet
    containing
    ``fig1`` /
    ``fig2`` /
    ``fig3``
    panels
    (very
    common
    in
    Nature /
    Science
    SI
    source
    data)
    is
    emitted
    as
    per-fig
    blocks
    with
    explicit
    row
    + column
    ranges
    so the
    LLM
    can
    tell
    panels
    apart
    and
    pass
    precise
    ranges
    to
    detectors.
    """
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ExtractionError(
            f"openpyxl not installed: {exc}"
        ) from exc
    try:
        # Note:
        # we
        # do
        # NOT
        # use
        # ``read_only=True``
        # because
        # the
        # fig
        # detector
        # needs
        # ``ws.max_row`` /
        # ``ws.max_column``
        # /
        # ``ws.cell(row, col).value``
        # which
        # are
        # not
        # available
        # on
        # ``ReadOnlyWorksheet``.
        wb = load_workbook(
            str(path), data_only=True
        )
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            f"xlsx open failed: {exc}"
        ) from exc
    try:
        out_lines: list[str] = []
        from ..ingest.xlsx import _cell_fill

        def render_cell(cell) -> str:
            value = cell.value
            text = "" if value is None else str(value)
            fill = _cell_fill(cell)
            if fill is None:
                return text
            return f"{text} [highlight:{fill}]"

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            out_lines.append(f"## Sheet: {sheet_name}")
            bboxes = detect_xlsx_figs(ws)
            if len(bboxes) <= 1:
                # No
                # multiple
                # figs
                # detected
                # --
                # emit
                # the
                # sheet
                # as
                # one
                # TSV
                # block
                # (preserves
                # the
                # Phase
                # B
                # behavior
                # for
                # the
                # common
                # single-table
                # case).
                for row in ws.iter_rows():
                    cells = [render_cell(cell) for cell in row]
                    out_lines.append("\t".join(cells))
            else:
                # Multiple
                # figs
                # detected
                # --
                # emit
                # one
                # fig
                # block
                # per
                # panel.
                # Use
                # a
                # per-fig
                # row
                # cap
                # so a
                # 20K-row
                # sheet
                # doesn't
                # produce
                # a
                # 100K-line
                # text
                # blob
                # in
                # the
                # LLM
                # context.
                # The
                # detector
                # reports
                # the
                # *full*
                # bbox
                # (rows
                # + cols)
                # so the
                # LLM
                # knows
                # exactly
                # what
                # was
                # truncated.
                for bb in bboxes:
                    n_rows = bb["bottom"] - bb["top"] + 1
                    n_cols = bb["right"] - bb["left"] + 1
                    out_lines.append(
                        f"  ## Fig: {bb['name']} "
                        f"(rows {bb['top']+1}-{bb['bottom']}, "
                        f"cols {bb['left']+1}-{bb['right']}, "
                        f"{n_rows} row(s) x {n_cols} col(s))"
                    )
                    _MAX_ROWS_PER_FIG = 200
                    if n_rows > _MAX_ROWS_PER_FIG:
                        out_lines.append(
                            f"    # showing first {_MAX_ROWS_PER_FIG} of {n_rows} rows; "
                            f"use offset to read the rest"
                        )
                    rendered = 0
                    for r in range(bb["top"], bb["bottom"] + 1):
                        if rendered >= _MAX_ROWS_PER_FIG:
                            out_lines.append(
                                f"    # ... ({n_rows - _MAX_ROWS_PER_FIG} more row(s) truncated)"
                            )
                            break
                        row_cells: list[str] = []
                        for c in range(
                            bb["left"], bb["right"] + 1
                        ):
                            cell = ws.cell(
                                row=r + 1, column=c + 1
                            )
                            row_cells.append(render_cell(cell))
                        out_lines.append("\t".join(row_cells))
                        rendered += 1
            out_lines.append("")
        wb.close()
        return "\n".join(out_lines)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            f"xlsx walk failed: {exc}"
        ) from exc


# =============================================================================
# Fig-boundary detector
# =============================================================================

# Match
# fig
# /
# table
# header
# cells
# like:
#   "Fig.S1a"
#   "Fig.S4b"
#   "Fig. S1a"
#   "Table S1"
#   "Tab.1"
#   "Fig. 2a"
# The
# regex
# is
# intentionally
# generous
# (any
# non-space
# word
# after
# Fig./Table)
# because
# the
# caller
# is
# downstream
# of
# openpyxl
# which
# may
# have
# a
# variety
# of
# spellings.
_FIG_HEADER_RE = re.compile(
    r"^\s*(?:Fig|Tab|Table)\.?\s*\S+",
    re.IGNORECASE,
)


def _looks_like_fig_header(cell_value: object) -> str | None:
    """Return the matched fig name if ``cell_value`` looks like a fig header.

    R-2026-06-17
    (Phase
    C):
    returns
    the
    full
    matched
    text
    (e.g.
    ``"Fig.S1a"``)
    so the
    caller
    can
    embed
    it in
    the
    output
    as a
    fig
    name.
    Returns
    ``None``
    for
    non-string
    or
    non-matching
    cells.
    """
    if not isinstance(cell_value, str):
        return None
    s = cell_value.strip()
    if not s:
        return None
    m = _FIG_HEADER_RE.match(s)
    if not m:
        return None
    # Strip
    # trailing
    # colon
    # /
    # whitespace
    # so the
    # fig
    # name
    # is
    # clean.
    return s.rstrip(":").strip()


def _is_blank_row(ws, row_idx: int, col_count: int) -> bool:
    """Return True if every cell in row ``row_idx`` (0-indexed) is None/empty.

    R-2026-06-17
    (Phase
    C):
    used
    to
    detect
    vertical
    splits
    (a
    fully
    empty
    row
    between
    two
    fig
    panels).
    """
    for c in range(col_count):
        v = ws.cell(row=row_idx + 1, column=c + 1).value
        if v is not None and str(v).strip():
            return False
    return True


def _is_blank_col(
    ws, col_idx: int, row_count: int, *, row_start: int = 0
) -> bool:
    """Return True if every cell in column ``col_idx`` (0-indexed) is None/empty.

    R-2026-06-17 (Phase C):
    used to detect horizontal
    splits (a fully empty
    column between two
    side-by-side fig panels).

    R-2026-06-19 (Phase D):
    added ``row_start`` so
    callers can scope the
    check to a specific
    row range. This matters
    for *vertical* figs: the
    cols to the right of a
    vertical fig are
    "filled" by *other* figs
    in the same sheet (the
    horizontal figs above)
    even though they are
    blank *for this fig's
    data*. Without
    ``row_start``, the
    detector would extend a
    vertical fig's right
    boundary all the way to
    the last data col of
    the horizontal figs.
    """
    for r in range(row_start, row_count):
        v = ws.cell(row=r + 1, column=col_idx + 1).value
        if v is not None and str(v).strip():
            return False
    return True


def detect_xlsx_figs(ws) -> list[dict[str, int | str]]:
    """Detect all fig/table panel boundaries in a single worksheet.

    R-2026-06-17
    (Phase
    C,
    fig-boundary
    detector):

    The
    returned
    list
    has
    one
    dict
    per
    panel
    with
    keys:

      * ``name`` --
        the
        matched
        fig
        header
        text
        (e.g.
        ``"Fig.S1a"``)
      * ``top`` /
        ``bottom`` --
        0-indexed
        row
        range
        (inclusive)
      * ``left`` /
        ``right`` --
        0-indexed
        col
        range
        (inclusive
        on
        left,
        exclusive
        on
        right
        for
        easy
        slicing
        via
        ``range(left, right)``)
      * ``header_row`` /
        ``header_col``
        --
        where
        the
        fig
        header
        was
        found

    Algorithm
    (two
    passes):

    1. **Header
    scan.**
    Look
    for
    fig
    headers
    (cells
    matching
    ``Fig./Table/...``)
    in
    the
    top
    2
    rows
    of
    the
    sheet
    (the
    most
    common
    position
    for
    fig
    titles).
    Then
    recursively
    walk
    down
    the
    sheet
    in
    50-row
    chunks
    to
    find
    more
    headers
    that
    might
    have
    been
    introduced
    by
    a
    vertical
    blank-row
    separator
    (e.g.
    Sfig.3
    Fig.S3a
    in
    R0
    +
    Fig.S3b
    in
    R8).

    2. **Bbox
    assignment.**
    For
    each
    header
    find
    the
    right
    + bottom
    boundary
    by
    (a)
    looking
    for
    the
    next
    header
    in
    the
    same
    row
    or
    same
    col,
    or
    (b)
    looking
    for
    the
    first
    blank
    row
    / col
    past
    the
    header,
    or
    (c)
    falling
    back
    to
    the
    sheet's
    max
    row
    / col.

    The
    output
    is
    *conservative*
    (boundaries
    include
    the
    next
    header
    row/col)
    so
    the
    LLM
    sees
    the
    full
    panel
    including
    its
    end-of-panel
    blank.
    """
    max_row = ws.max_row
    max_col = ws.max_column
    if max_row == 0 or max_col == 0:
        return []

    headers: list[tuple[str, int, int]] = []

    def scan_band(r_start: int, r_end: int) -> None:
        for r in range(r_start, min(r_end, max_row)):
            for c in range(max_col):
                name = _looks_like_fig_header(
                    ws.cell(row=r + 1, column=c + 1).value
                )
                if name:
                    headers.append((name, r, c))

    # Pass 1: top header band.
    scan_band(0, 2)
    # Pass 2: walk down to find headers below the first band
    # (e.g. Sfig.3 Fig.S3b at R8, Sfig.4 Fig.S4b at R6).
    last_seen = max((h[1] for h in headers), default=-1)
    chunk = 50
    r = last_seen + 1
    while r < max_row:
        scan_band(r, r + chunk)
        new = [h for h in headers if h[1] >= r]
        if new:
            r = max(h[1] for h in new) + 1
        else:
            r += chunk

    # Dedupe by name (same fig may be detected multiple times if
    # header is in both the top band and a later band).
    seen: set[str] = set()
    unique: list[tuple[str, int, int]] = []
    for h in headers:
        if h[0] in seen:
            continue
        seen.add(h[0])
        unique.append(h)
    if not unique:
        return []

    # Group headers by row (for horizontal-split neighbors) and by col
    # (for vertical-split neighbors).
    by_row: dict[int, list[tuple[str, int, int]]] = {}
    by_col: dict[int, list[tuple[str, int, int]]] = {}
    for h in unique:
        by_row.setdefault(h[1], []).append(h)
        by_col.setdefault(h[2], []).append(h)
    for r in by_row:
        by_row[r].sort(key=lambda x: x[2])
    for c in by_col:
        by_col[c].sort(key=lambda x: x[1])

    bboxes: list[dict[str, int | str]] = []
    for name, hdr_row, hdr_col in unique:
        # Right boundary.
        # R-2026-06-19 (Phase D):
        # the fig-boundary
        # detector uses an
        # *exclusive* right
        # boundary (so
        # ``range(left, right)``
        # yields the cols of
        # this fig). For
        # storage in
        # ``ExtractedTable.bbox``
        # we convert to
        # *inclusive* right
        # (``right - 1``) so
        # the 1-indexed
        # ``ListDataSourcesTool``
        # output can show
        # ``cols 1-3`` without
        # off-by-one confusion.
        right_excl: int | None = None
        for h2 in by_row[hdr_row]:
            if h2[2] > hdr_col:
                right_excl = h2[2]
                break
        if right_excl is None:
            c = hdr_col + 1
            while c < max_col:
                # R-2026-06-19 (Phase D):
                # scope the blank-col
                # check to the fig's
                # own row range (from
                # the header row down)
                # so that data in
                # *other* figs in
                # the same sheet
                # doesn't fool us
                # into extending the
                # right boundary to
                # the last data col
                # of those other figs.
                if _is_blank_col(
                    ws, c, max_row,
                    row_start=hdr_row,
                ):
                    right_excl = c
                    break
                c += 1
        if right_excl is None:
            right_excl = max_col

        # Bottom boundary (exclusive).
        bottom_excl: int | None = None
        for h2 in by_col[hdr_col]:
            if h2[1] > hdr_row:
                bottom_excl = h2[1]
                break
        if bottom_excl is None:
            r = hdr_row + 1
            while r < max_row:
                if _is_blank_row(ws, r, max_col):
                    bottom_excl = r
                    break
                r += 1
        if bottom_excl is None:
            bottom_excl = max_row

        bboxes.append({
            "name": name,
            "top": hdr_row,
            "left": hdr_col,
            # Inclusive right (last col of THIS fig).
            "right": right_excl - 1,
            # Inclusive bottom (last row of THIS fig).
            "bottom": bottom_excl - 1,
            "header_row": hdr_row,
            "header_col": hdr_col,
        })
    return bboxes


def extract_pptx_text(path: str) -> str:
    """Extract plain text from a .pptx (PowerPoint) file."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ExtractionError(
            f"python-pptx not installed: {exc}"
        ) from exc
    try:
        pres = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            f"pptx open failed: {exc}"
        ) from exc
    try:
        out_lines: list[str] = []
        for i, slide in enumerate(pres.slides, start=1):
            out_lines.append(f"## Slide {i}")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = "".join(
                        run.text for run in para.runs
                    ).strip()
                    if text:
                        out_lines.append(text)
            if slide.has_notes_slide:
                notes = (
                    slide.notes_slide.notes_text_frame.text or ""
                ).strip()
                if notes:
                    out_lines.append(f"[notes] {notes}")
            out_lines.append("")
        return "\n".join(out_lines)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            f"pptx walk failed: {exc}"
        ) from exc


def extract_ipynb_text(path: str) -> str:
    """Extract a plain-text rendering of a .ipynb (Jupyter) notebook."""
    try:
        import nbformat
    except ImportError as exc:
        raise ExtractionError(
            f"nbformat not installed: {exc}"
        ) from exc
    try:
        nb = nbformat.read(str(path), as_version=4)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            f"ipynb parse failed: {exc}"
        ) from exc
    try:
        out_lines: list[str] = []
        for i, cell in enumerate(nb.cells, start=1):
            cell_type = cell.get("cell_type", "")
            source = cell.get("source", "")
            if isinstance(source, list):
                source = "".join(source)
            if cell_type == "markdown":
                out_lines.append(f"# Cell {i} (markdown)")
                out_lines.append(source.rstrip())
            elif cell_type == "code":
                lang = (
                    nb.metadata.get("kernelspec", {}).get(
                        "language", "python"
                    )
                )
                out_lines.append(f"# Cell {i} (code, {lang})")
                out_lines.append(source.rstrip())
                outputs = cell.get("outputs", [])
                for out in outputs:
                    otype = out.get("output_type", "")
                    if otype == "stream":
                        text = out.get("text", "")
                        if isinstance(text, list):
                            text = "".join(text)
                        if text:
                            out_lines.append(
                                f"[stdout] {text.rstrip()[:500]}"
                            )
                    elif otype in (
                        "execute_result",
                        "display_data",
                    ):
                        data = out.get("data", {})
                        if "text/plain" in data:
                            tp = data["text/plain"]
                            if isinstance(tp, list):
                                tp = "".join(tp)
                            out_lines.append(
                                f"[result] {tp.rstrip()[:500]}"
                            )
                    elif otype == "error":
                        ename = out.get("ename", "")
                        evalue = out.get("evalue", "")
                        out_lines.append(
                            f"[error] {ename}: {evalue}"
                        )
            out_lines.append("")
        return "\n".join(out_lines)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            f"ipynb walk failed: {exc}"
        ) from exc


_EXTRACTORS = {
    ".docx": extract_docx_text,
    ".xlsx": extract_xlsx_text,
    ".pptx": extract_pptx_text,
    ".ipynb": extract_ipynb_text,
}


def try_extract_document_real(
    path: str,
    *,
    on_error: str = "fallback",
) -> str | None:
    """Real implementation of ``try_extract_document`` (Phase A #8).

    R-2026-06-17
    (Phase B):
    same
    shape
    as
    the
    Phase
    A
    stub
    in
    ``safe_read.py``
    (signature
    ``(path, *, on_error="fallback") -> str | None``)
    but
    actually
    walks
    the
    document
    structure.
    """
    if not path or not Path(path).exists():
        return None
    suffix = Path(path).suffix.lower()
    extract_fn = _EXTRACTORS.get(suffix)
    if extract_fn is None:
        return None
    try:
        return extract_fn(path)
    except ExtractionError:
        if on_error == "fallback":
            return None
        raise
    except Exception as exc:  # noqa: BLE001
        if on_error == "fallback":
            return None
        raise ExtractionError(
            f"{suffix} extraction failed: {exc}"
        ) from exc

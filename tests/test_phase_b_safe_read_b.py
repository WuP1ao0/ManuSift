"""R-2026-06-17 (Phase B:
borrow from
Claude Code +
Hermes):
4 medium-cost
``read_file``
hardening
modules in
``manusift/tools/safe_read_b.py``,
all wired into
``ReadFileTool.execute``
(``manusift/tools/direct_fs.py``).

The four
guards are:

  * B-1: ``suggest_similar_files``
    — fuzzy
    match
    when
    a
    file
    is
    not
    found
  * B-2: ``ReadTracker``
    — mtime
    dedup
    + BLOCKED-after-2-hits
  * B-3: ``redact_sensitive_text``
    — replace
    30+
    known
    API-key
    prefixes
    with
    ``[REDACTED]``
  * B-4: real
    ``python-docx``
    /
    ``openpyxl``
    /
    ``python-pptx``
    /
    ``nbformat``
    extractors

The tests
below
assert
each
guard
fires
in
isolation
(B-1..B-4
unit
tests)
and
in
combination
(end-to-end
ReadFileTool
tests).
"""
from __future__ import annotations

import json
import os
import sys
import tempfile
import time
import zipfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from manusift.tools.safe_read_b import (
    _READ_DEDUP_STATUS_MESSAGE,
    _score_similarity,
    ReadTracker,
    extract_docx_text,
    extract_ipynb_text,
    extract_pptx_text,
    extract_xlsx_text,
    redact_sensitive_text,
    suggest_similar_files,
    try_extract_document_real,
    ExtractionError,
)
from manusift.tools.direct_fs import ReadFileTool
from manusift.tools.tool import ToolContext


@pytest.fixture(autouse=True)
def _reset_trackers():
    """Reset the per-trace ReadTracker cache between tests so dedup state doesn't leak."""
    from manusift.tools.safe_read_b import reset_all_trackers
    reset_all_trackers()
    yield
    reset_all_trackers()


@pytest.fixture
def workspace(tmp_path, monkeypatch):
    """Redirect MANUSIFT_WORKSPACE_DIR so direct_fs is happy."""
    monkeypatch.setenv("MANUSIFT_WORKSPACE_DIR", str(tmp_path))
    yield tmp_path


@pytest.fixture
def ctx():
    return ToolContext(trace_id="test_trace", current_pdf=None, metadata={})


@pytest.fixture
def read_file_tool():
    return ReadFileTool()


# ============================================================================
# B-1: suggest_similar_files (unit tests)
# ============================================================================


class TestScoreSimilarity:
    @pytest.mark.parametrize(
        "q,c,ext,expected",
        [
            ("notes.md", "notes.md", ".md", 100),  # exact
            ("notes.md", "notes.txt", ".txt", 90),  # same base, diff ext
            ("notes.md", "notes_old.md", ".md", 70),  # prefix
            ("notes.md", "my_notes.md", ".md", 60),  # substring
            ("notes.md", "votes.md", ".md", 30),  # same ext + 40% overlap
            ("notes.md", "image.png", ".png", 0),  # no match
        ],
    )
    def test_known_scores(self, q, c, ext, expected):
        assert _score_similarity(q, c, ext) == expected

    def test_case_insensitive(self):
        assert _score_similarity("Notes.MD", "notes.md", ".md") == 100

    def test_empty_inputs(self):
        assert _score_similarity("", "foo.md", ".md") == 0
        assert _score_similarity("foo.md", "", ".md") == 0


class TestSuggestSimilarFiles:
    def test_returns_top_match(self, tmp_path):
        (tmp_path / "notes.md").write_text("hi", encoding="utf-8")
        (tmp_path / "other.md").write_text("bye", encoding="utf-8")
        result = suggest_similar_files(
            str(tmp_path / "notes.md"),
            search_dir=str(tmp_path),
        )
        # ``notes.md``
        # is the
        # top
        # match
        # (score=100,
        # vs
        # ``other.md``=0).
        assert len(result) >= 1
        assert result[0].endswith("notes.md")

    def test_returns_extension_mismatch_first(
        self, tmp_path
    ):
        (tmp_path / "notes.txt").write_text("hi", encoding="utf-8")
        (tmp_path / "random.md").write_text("bye", encoding="utf-8")
        # User
        # asked
        # for
        # ``notes.md``
        # but
        # only
        # ``notes.txt``
        # exists
        # (same
        # base,
        # different
        # ext
        # →
        # score=90).
        result = suggest_similar_files(
            str(tmp_path / "notes.md"),
            search_dir=str(tmp_path),
        )
        assert len(result) >= 1
        assert result[0].endswith("notes.txt")

    def test_returns_absolute_paths(self, tmp_path):
        (tmp_path / "notes.md").write_text("hi", encoding="utf-8")
        result = suggest_similar_files(
            str(tmp_path / "notes.md"),
            search_dir=str(tmp_path),
        )
        for p in result:
            assert os.path.isabs(p)

    def test_empty_dir(self, tmp_path):
        result = suggest_similar_files(
            str(tmp_path / "nope.md"),
            search_dir=str(tmp_path),
        )
        assert result == []

    def test_nonexistent_search_dir(self, tmp_path):
        result = suggest_similar_files(
            str(tmp_path / "nope.md"),
            search_dir=str(tmp_path / "no_such_dir"),
        )
        assert result == []

    def test_limit_respected(self, tmp_path):
        # Create 10
        # files
        # all
        # with
        # similar
        # names.
        for i in range(10):
            (tmp_path / f"notes_{i}.md").write_text("x", encoding="utf-8")
        result = suggest_similar_files(
            str(tmp_path / "notes.md"),
            search_dir=str(tmp_path),
            limit=3,
        )
        assert len(result) <= 3

    def test_skips_subdirectories(self, tmp_path):
        sub = tmp_path / "sub"
        sub.mkdir()
        (sub / "notes.md").write_text("hi", encoding="utf-8")
        # ``suggest_similar_files``
        # is
        # non-recursive
        # by
        # design
        # -- only
        # the
        # top-level
        # dir
        # is
        # scanned.
        # (The
        # LLM
        # can
        # call
        # ``list_dir``
        # if it
        # needs
        # recursive
        # suggestions.)
        result = suggest_similar_files(
            str(tmp_path / "notes.md"),
            search_dir=str(tmp_path),
        )
        assert result == []


# ============================================================================
# B-2: ReadTracker (unit tests)
# ============================================================================


class TestReadTracker:
    def test_first_call_proceeds(self):
        t = ReadTracker()
        result = t.check(
            "foo.md", 1, 100, current_mtime=100.0
        )
        assert result is None  # proceed with real read

    def test_second_call_returns_unchanged_stub(self):
        t = ReadTracker()
        t.record("foo.md", 1, 100, mtime=100.0)
        result = t.check(
            "foo.md", 1, 100, current_mtime=100.0
        )
        assert result is not None
        assert result["status"] == "unchanged"
        assert result["dedup"] is True
        assert result["content_returned"] is False
        assert "message" in result
        assert _READ_DEDUP_STATUS_MESSAGE in result["message"]

    def test_third_call_returns_blocked(self):
        t = ReadTracker()
        t.record("foo.md", 1, 100, mtime=100.0)
        t.check("foo.md", 1, 100, current_mtime=100.0)  # 1st hit
        result = t.check(
            "foo.md", 1, 100, current_mtime=100.0
        )  # 2nd hit
        assert result is not None
        assert result["ok"] is False
        assert result["error_kind"] == "blocked"
        assert "BLOCKED" in result["error"]
        assert "STOP" in result["error"]
        assert result["already_read"] >= 2

    def test_file_changed_proceeds(self):
        t = ReadTracker()
        t.record("foo.md", 1, 100, mtime=100.0)
        # mtime
        # changed
        # →
        # proceed.
        result = t.check(
            "foo.md", 1, 100, current_mtime=200.0
        )
        assert result is None

    def test_different_offset_proceeds(self):
        t = ReadTracker()
        t.record("foo.md", 1, 100, mtime=100.0)
        # Different
        # (offset,
        # limit)
        # →
        # proceed.
        result = t.check(
            "foo.md", 200, 100, current_mtime=100.0
        )
        assert result is None

    def test_different_path_proceeds(self):
        t = ReadTracker()
        t.record("foo.md", 1, 100, mtime=100.0)
        result = t.check(
            "bar.md", 1, 100, current_mtime=100.0
        )
        assert result is None

    def test_reset(self):
        t = ReadTracker()
        t.record("foo.md", 1, 100, mtime=100.0)
        t.reset()
        result = t.check(
            "foo.md", 1, 100, current_mtime=100.0
        )
        assert result is None  # tracker cleared

    def test_cap_enforced(self):
        t = ReadTracker()
        # Force
        # the
        # cap
        # by
        # recording
        # 1001
        # distinct
        # files
        # (read
        # tracker
        # caps
        # at
        # ``_DEDUP_CAP = 1000``).
        for i in range(1001):
            t.record(f"foo_{i}.md", 1, 100, mtime=100.0)
        # After
        # eviction,
        # oldest
        # entries
        # (foo_0.md
        # ..
        # foo_0.md)
        # are
        # gone.
        # ``foo_0.md``
        # should
        # not be
        # a
        # dedup
        # hit
        # anymore
        # (i.e.
        # the
        # tracker
        # has
        # forgotten
        # it).
        result = t.check("foo_0.md", 1, 100, current_mtime=100.0)
        # Either
        # ``None``
        # (forgotten
        # → proceed)
        # or
        # "unchanged"
        # stub
        # if
        # the
        # cap
        # evicted
        # some
        # other
        # entry
        # instead.
        # Both
        # are
        # acceptable
        # -- the
        # important
        # property
        # is that
        # the
        # dict
        # size
        # is
        # capped.
        # Check
        # the
        # dict
        # size
        # directly:
        assert len(t._dedup) <= 1000


# ============================================================================
# B-3: redact_sensitive_text (unit tests)
# ============================================================================


class TestRedactSensitiveText:
    @pytest.mark.parametrize(
        "secret",
        [
            "sk-1234567890abcdefghij",
            "sk-proj-1234567890abcdefghij",
            "ghp_abcdefghij1234567890",
            "github_pat_11ABCDEFG0_abc123def456ghi789jkl012mno345pqr",
            "gho_abcdefghij1234567890",
            "xoxb-1234567890-abcdefghij",
            "AIzaSyAbcdefghijklmnopqrstuvwxyz1234567",
            "AKIAIOSFODNN7EXAMPLE",
            "sk_live_abcdefghij1234567890",
            "sk_test_abcdefghij1234567890",
            "SG.abcdefghij1234567890.abcdefghij",
            "hf_abcdefghij1234567890",
            "r8_abcdefghij1234567890",
            "npm_abcdefghij1234567890",
            "pypi-abcdefghij1234567890",
            "gAAAAAabcdefghij1234567890_=-_",
        ],
    )
    def test_known_secrets_redacted(self, secret):
        text = f"api_key = {secret}"
        out = redact_sensitive_text(text)
        assert secret not in out
        assert "[REDACTED]" in out

    def test_surrounding_code_preserved(self):
        text = "client = OpenAI(api_key=sk-abcdefghijklmnopqrst)"
        out = redact_sensitive_text(text)
        assert "client = OpenAI(" in out
        assert "api_key=" in out

    def test_placeholder_kwarg(self):
        text = "key = sk-abcdefghijklmnopqrst"
        out = redact_sensitive_text(text, placeholder="[HIDDEN]")
        assert "[HIDDEN]" in out
        assert "[REDACTED]" not in out

    def test_empty(self):
        assert redact_sensitive_text("") == ""

    def test_no_secrets_unchanged(self):
        text = "def hello():\n    print('hi')\n"
        out = redact_sensitive_text(text)
        assert out == text

    def test_word_boundary_protects_substrings(self):
        # ``skull``
        # should
        # NOT
        # be
        # matched
        # as
        # ``sk-``
        # +
        # ``ull``
        # because
        # the
        # boundary
        # check
        # requires
        # non-alnum
        # around
        # the
        # match.
        text = "the word skull is not a secret"
        out = redact_sensitive_text(text)
        assert "skull" in out
        assert "[REDACTED]" not in out

    def test_minimum_length_enforced(self):
        # The
        # regex
        # requires
        # the
        # prefix
        # +
        # at
        # least
        # N
        # token
        # chars.
        # A
        # short
        # token
        # like
        # ``sk-abc``
        # (only
        # 3
        # chars)
        # is
        # not
        # redacted.
        text = "short = sk-abc"
        out = redact_sensitive_text(text)
        assert "sk-abc" in out


# ============================================================================
# B-4: structured document extractors (unit tests)
# ============================================================================


class TestExtractDocx:
    def test_minimal_docx(self, tmp_path):
        # Create
        # a
        # minimal
        # valid
        # .docx
        # file
        # using
        # ``python-docx``.
        from docx import Document
        doc = Document()
        doc.add_paragraph("Hello, world!")
        doc.add_paragraph("Second paragraph.")
        path = tmp_path / "test.docx"
        doc.save(str(path))
        text = extract_docx_text(str(path))
        assert "Hello, world!" in text
        assert "Second paragraph." in text

    def test_corrupted_docx_raises_extraction_error(
        self, tmp_path
    ):
        bad = tmp_path / "bad.docx"
        bad.write_bytes(b"not a real docx")
        with pytest.raises(ExtractionError):
            extract_docx_text(str(bad))

    def test_nonexistent_file_raises(self, tmp_path):
        with pytest.raises(ExtractionError):
            extract_docx_text(str(tmp_path / "nope.docx"))


class TestExtractXlsx:
    def test_minimal_xlsx(self, tmp_path):
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["Name", "Age"])
        ws.append(["Alice", 30])
        path = tmp_path / "test.xlsx"
        wb.save(str(path))
        text = extract_xlsx_text(str(path))
        assert "## Sheet: Sheet1" in text
        assert "Name" in text
        assert "Age" in text
        assert "Alice" in text
        assert "30" in text

    def test_xlsx_highlight_marker_is_rendered(self, tmp_path):
        from openpyxl import Workbook
        from openpyxl.styles import PatternFill
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["Name", "Value"])
        ws.append(["Alice", 22.2])
        ws["B2"].fill = PatternFill(
            fill_type="solid",
            fgColor="FFFF00",
        )
        path = tmp_path / "highlighted.xlsx"
        wb.save(str(path))

        text = extract_xlsx_text(str(path))

        assert "22.2 [highlight:FFFF00]" in text

    def test_multiple_sheets(self, tmp_path):
        from openpyxl import Workbook
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "First"
        ws1.append(["A"])
        ws2 = wb.create_sheet("Second")
        ws2.append(["B"])
        path = tmp_path / "multi.xlsx"
        wb.save(str(path))
        text = extract_xlsx_text(str(path))
        assert "## Sheet: First" in text
        assert "## Sheet: Second" in text

    def test_corrupted_xlsx_raises(self, tmp_path):
        bad = tmp_path / "bad.xlsx"
        bad.write_bytes(b"not a real xlsx")
        with pytest.raises(ExtractionError):
            extract_xlsx_text(str(bad))


class TestExtractPptx:
    def test_minimal_pptx(self, tmp_path):
        from pptx import Presentation
        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[0])
        slide.shapes.title.text = "Title Slide"
        path = tmp_path / "test.pptx"
        pres.save(str(path))
        text = extract_pptx_text(str(path))
        assert "## Slide 1" in text
        assert "Title Slide" in text

    def test_corrupted_pptx_raises(self, tmp_path):
        bad = tmp_path / "bad.pptx"
        bad.write_bytes(b"not a real pptx")
        with pytest.raises(ExtractionError):
            extract_pptx_text(str(bad))


class TestExtractIpynb:
    def test_minimal_ipynb(self, tmp_path):
        # A
        # valid
        # .ipynb
        # is
        # a
        # JSON
        # file
        # with
        # the
        # ``nbformat``
        # schema.
        nb = {
            "cells": [
                {
                    "cell_type": "markdown",
                    "metadata": {},
                    "source": ["# Heading\n", "Some text"],
                },
                {
                    "cell_type": "code",
                    "execution_count": 1,
                    "metadata": {},
                    "source": ["print('hello')\n"],
                    "outputs": [
                        {
                            "output_type": "stream",
                            "name": "stdout",
                            "text": ["hello\n"],
                        }
                    ],
                },
            ],
            "metadata": {
                "kernelspec": {
                    "name": "python3",
                    "language": "python",
                }
            },
            "nbformat": 4,
            "nbformat_minor": 5,
        }
        import json
        path = tmp_path / "test.ipynb"
        path.write_text(json.dumps(nb), encoding="utf-8")
        text = extract_ipynb_text(str(path))
        assert "# Cell 1 (markdown)" in text
        assert "# Heading" in text
        assert "# Cell 2 (code, python)" in text
        assert "print('hello')" in text
        assert "[stdout] hello" in text

    def test_corrupted_ipynb_raises(self, tmp_path):
        bad = tmp_path / "bad.ipynb"
        bad.write_text("not valid JSON", encoding="utf-8")
        with pytest.raises(ExtractionError):
            extract_ipynb_text(str(bad))


class TestTryExtractDocumentReal:
    def test_returns_text_for_valid_docx(self, tmp_path):
        from docx import Document
        doc = Document()
        doc.add_paragraph("test paragraph")
        path = tmp_path / "test.docx"
        doc.save(str(path))
        text = try_extract_document_real(str(path))
        assert text is not None
        assert "test paragraph" in text

    def test_returns_none_for_corrupted_docx(
        self, tmp_path
    ):
        bad = tmp_path / "bad.docx"
        bad.write_bytes(b"not a real docx")
        text = try_extract_document_real(
            str(bad), on_error="fallback"
        )
        assert text is None

    def test_returns_none_for_unsupported_ext(
        self, tmp_path
    ):
        f = tmp_path / "test.txt"
        f.write_text("hello", encoding="utf-8")
        text = try_extract_document_real(str(f))
        assert text is None

    def test_returns_none_for_missing_file(
        self, tmp_path
    ):
        text = try_extract_document_real(
            str(tmp_path / "nope.docx")
        )
        assert text is None


# ============================================================================
# End-to-end: ReadFileTool.execute with B-1..B-4
# ============================================================================


class TestReadFileToolB1SimilarFiles:
    def test_not_found_returns_suggestions(
        self, read_file_tool, ctx, tmp_path
    ):
        # Create
        # a
        # similar-named
        # file
        # in the
        # same
        # dir.
        (tmp_path / "notes.md").write_text(
            "hi", encoding="utf-8"
        )
        out = read_file_tool.execute(
            {"path": str(tmp_path / "notes.txt")}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is False
        assert data["error_kind"] == "not_found"
        assert "similar_files" in data
        assert any(
            p.endswith("notes.md")
            for p in data["similar_files"]
        )
        assert "Did you mean" in data["error"]

    def test_not_found_no_similar(
        self, read_file_tool, ctx, tmp_path
    ):
        out = read_file_tool.execute(
            {"path": str(tmp_path / "absent.md")}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is False
        assert "similar_files" not in data


class TestReadFileToolB2Dedup:
    def test_first_read_succeeds(
        self, read_file_tool, ctx, tmp_path
    ):
        f = tmp_path / "f.md"
        f.write_text("hello world", encoding="utf-8")
        out = read_file_tool.execute(
            {"path": str(f)}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is True
        assert "hello world" in data["content"]

    def test_second_read_unchanged_stub(
        self, read_file_tool, ctx, tmp_path
    ):
        f = tmp_path / "f.md"
        f.write_text("hello world", encoding="utf-8")
        # First
        # read:
        # real
        # content.
        read_file_tool.execute({"path": str(f)}, ctx)
        # Second
        # read
        # (file
        # unchanged):
        # stub.
        out = read_file_tool.execute(
            {"path": str(f)}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is True
        assert data["status"] == "unchanged"
        assert data["content_returned"] is False

    def test_third_read_blocked(
        self, read_file_tool, ctx, tmp_path
    ):
        f = tmp_path / "f.md"
        f.write_text("hello", encoding="utf-8")
        read_file_tool.execute({"path": str(f)}, ctx)
        read_file_tool.execute({"path": str(f)}, ctx)
        out = read_file_tool.execute(
            {"path": str(f)}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is False
        assert data["error_kind"] == "blocked"
        assert "BLOCKED" in data["error"]

    def test_file_modified_proceeds(
        self, read_file_tool, ctx, tmp_path
    ):
        f = tmp_path / "f.md"
        f.write_text("v1", encoding="utf-8")
        read_file_tool.execute({"path": str(f)}, ctx)
        # Modify
        # the
        # file
        # →
        # mtime
        # changes
        # →
        # next
        # read
        # returns
        # real
        # content.
        time.sleep(0.05)
        f.write_text("v2", encoding="utf-8")
        out = read_file_tool.execute(
            {"path": str(f)}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is True
        assert "v2" in data["content"]

    def test_different_session_no_dedup(
        self, read_file_tool, workspace, tmp_path
    ):
        f = tmp_path / "f.md"
        f.write_text("hello", encoding="utf-8")
        # Session 1
        ctx1 = ToolContext(
            trace_id="s1", current_pdf=None, metadata={}
        )
        read_file_tool.execute({"path": str(f)}, ctx1)
        # Session 2
        # (different
        # trace_id
        # → different
        # ctx.metadata
        # → no
        # dedup)
        ctx2 = ToolContext(
            trace_id="s2", current_pdf=None, metadata={}
        )
        out = read_file_tool.execute(
            {"path": str(f)}, ctx2
        )
        data = json.loads(out)
        assert data["ok"] is True
        # Real
        # content
        # (not
        # stub).
        assert "status" not in data or data.get("status") != "unchanged"


class TestReadFileToolB3Redaction:
    def test_secret_redacted(
        self, read_file_tool, ctx, tmp_path
    ):
        f = tmp_path / "config.py"
        f.write_text(
            "OPENAI_KEY = 'sk-abcdefghijklmnopqrst'\n",
            encoding="utf-8",
        )
        out = read_file_tool.execute(
            {"path": str(f)}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is True
        # The
        # key
        # is
        # gone.
        assert "sk-abcdefghijklmnopqrst" not in data["content"]
        # The
        # surrounding
        # code
        # is
        # preserved.
        assert "OPENAI_KEY" in data["content"]

    def test_no_secret_unchanged(
        self, read_file_tool, ctx, tmp_path
    ):
        f = tmp_path / "notes.md"
        f.write_text("just a note", encoding="utf-8")
        out = read_file_tool.execute(
            {"path": str(f)}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is True
        assert "just a note" in data["content"]


class TestReadFileToolB4DocxExtract:
    def test_docx_returns_extracted_text(
        self, read_file_tool, ctx, tmp_path
    ):
        from docx import Document
        doc = Document()
        doc.add_paragraph("Hello, DOCX world!")
        path = tmp_path / "real.docx"
        doc.save(str(path))
        out = read_file_tool.execute(
            {"path": str(path)}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is True
        assert data["extracted_document"] is True
        assert "Hello, DOCX world!" in data["content"]

    def test_corrupted_docx_falls_through(
        self, read_file_tool, ctx, tmp_path
    ):
        bad = tmp_path / "bad.docx"
        bad.write_bytes(b"not a real docx")
        out = read_file_tool.execute(
            {"path": str(bad)}, ctx
        )
        data = json.loads(out)
        # Extraction
        # failed
        # (ExtractionError)
        # → fallback
        # to
        # text-read
        # →
        # either
        # succeeds
        # (latin-1
        # decode
        # of
        # the
        # raw
        # bytes)
        # or
        # fails
        # with
        # an
        # error
        # JSON.
        # Either
        # way,
        # no
        # crash.
        assert isinstance(data, dict)
        assert "ok" in data

    def test_xlsx_returns_extracted_text(
        self, read_file_tool, ctx, tmp_path
    ):
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.append(["Header"])
        ws.append(["value1"])
        path = tmp_path / "real.xlsx"
        wb.save(str(path))
        out = read_file_tool.execute(
            {"path": str(path)}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is True
        assert data["extracted_document"] is True
        assert "## Sheet:" in data["content"]
        assert "value1" in data["content"]


# ============================================================================
# Regression: existing behavior preserved
# ============================================================================


class TestReadFileToolRegression:
    def test_normal_text_file(
        self, read_file_tool, ctx, tmp_path
    ):
        f = tmp_path / "notes.md"
        f.write_text("line 1\nline 2\n", encoding="utf-8")
        out = read_file_tool.execute({"path": str(f)}, ctx)
        data = json.loads(out)
        assert data["ok"] is True
        assert "line 1" in data["content"]
        assert "line 2" in data["content"]

    def test_pdf_still_rejected(
        self, read_file_tool, ctx, tmp_path
    ):
        f = tmp_path / "paper.pdf"
        f.write_bytes(b"%PDF-1.4 fake")
        out = read_file_tool.execute({"path": str(f)}, ctx)
        data = json.loads(out)
        assert data["ok"] is False
        assert "ingest_from_path" in data["error"]

    def test_png_still_rejected(
        self, read_file_tool, ctx, tmp_path
    ):
        f = tmp_path / "image.png"
        f.write_bytes(b"\x89PNG fake")
        out = read_file_tool.execute({"path": str(f)}, ctx)
        data = json.loads(out)
        assert data["ok"] is False
        assert "binary" in data["error"].lower()

    def test_dev_zero_still_blocked(
        self, read_file_tool, ctx
    ):
        out = read_file_tool.execute(
            {"path": "/dev/zero"}, ctx
        )
        data = json.loads(out)
        assert data["ok"] is False
        assert "device file" in data["error"]
